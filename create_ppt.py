from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import os

def create_presentation():
    # Create presentation
    prs = Presentation()
    
    # Define paths to screenshots (using the ones we verified exist)
    base_path = r"C:\Users\Imamsab jamdar\.gemini\antigravity\brain\191c3fa0-c7a0-4d2a-b173-bc27da7a4fc1"
    dashboard_img = os.path.join(base_path, "dashboard_main_view_1764654012628.png")
    predictions_img = os.path.join(base_path, "predictions_page_view_1764654066655.png")
    
    # Helper to add slide with title and content
    def add_slide(title, content_text=None):
        slide_layout = prs.slide_layouts[1] # Title and Content
        slide = prs.slides.add_slide(slide_layout)
        slide.shapes.title.text = title
        if content_text:
            tf = slide.shapes.placeholders[1].text_frame
            tf.text = content_text
        return slide

    # 1. Title Slide
    slide_layout = prs.slide_layouts[0] # Title Slide
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "Predictive Maintenance Dashboard\nfor Solar Panels"
    subtitle.text = "Major Project Presentation\nDecember 2025"

    # 2. Introduction
    slide = add_slide("Introduction")
    tf = slide.shapes.placeholders[1].text_frame
    tf.text = "Project Overview"
    p = tf.add_paragraph()
    p.text = "• An intelligent monitoring system for solar panel arrays"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• Combines real-time data monitoring with machine learning"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• Predicts maintenance needs to prevent downtime"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• Optimizes energy output and extends equipment lifespan"
    p.level = 1

    # 3. Technology Stack
    slide = add_slide("Technology Stack")
    tf = slide.shapes.placeholders[1].text_frame
    tf.text = "Frontend"
    p = tf.add_paragraph()
    p.text = "• React.js 19.0, Tailwind CSS, Recharts"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "Backend"
    p = tf.add_paragraph()
    p.text = "• Python FastAPI, Uvicorn"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "Database & ML"
    p = tf.add_paragraph()
    p.text = "• MongoDB (Data Storage)"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• Scikit-learn (Random Forest Models)"
    p.level = 1

    # 4. Dashboard Overview (with Screenshot)
    slide = prs.slides.add_slide(prs.slide_layouts[5]) # Title Only
    slide.shapes.title.text = "Dashboard Overview"
    
    if os.path.exists(dashboard_img):
        left = Inches(0.5)
        top = Inches(1.5)
        height = Inches(5.5)
        slide.shapes.add_picture(dashboard_img, left, top, height=height)
        
        # Add some text annotations
        txBox = slide.shapes.add_textbox(Inches(6.5), Inches(2), Inches(3), Inches(3))
        tf = txBox.text_frame
        tf.text = "Key Features:"
        p = tf.add_paragraph()
        p.text = "• Real-time KPIs (Power, Efficiency)"
        p.level = 0
        p = tf.add_paragraph()
        p.text = "• Live Performance Charts"
        p.level = 0
        p = tf.add_paragraph()
        p.text = "• Maintenance Status Indicators"
        p.level = 0
    else:
        txBox = slide.shapes.add_textbox(Inches(2), Inches(3), Inches(6), Inches(2))
        txBox.text_frame.text = "[Dashboard Screenshot Placeholder]"

    # 5. Predictive Analytics (with Screenshot)
    slide = prs.slides.add_slide(prs.slide_layouts[5]) # Title Only
    slide.shapes.title.text = "AI-Driven Predictions"
    
    if os.path.exists(predictions_img):
        left = Inches(0.5)
        top = Inches(1.5)
        height = Inches(5.5)
        slide.shapes.add_picture(predictions_img, left, top, height=height)
        
        txBox = slide.shapes.add_textbox(Inches(6.5), Inches(2), Inches(3), Inches(3))
        tf = txBox.text_frame
        tf.text = "ML Capabilities:"
        p = tf.add_paragraph()
        p.text = "• 7-Day Power Forecast"
        p.level = 0
        p = tf.add_paragraph()
        p.text = "• Efficiency Degradation Prediction"
        p.level = 0
        p = tf.add_paragraph()
        p.text = "• Maintenance Recommendations"
        p.level = 0
    else:
        txBox = slide.shapes.add_textbox(Inches(2), Inches(3), Inches(6), Inches(2))
        txBox.text_frame.text = "[Predictions Screenshot Placeholder]"

    # 6. Key Features
    slide = add_slide("Key Features")
    tf = slide.shapes.placeholders[1].text_frame
    tf.text = "Comprehensive Monitoring"
    p = tf.add_paragraph()
    p.text = "• Real-time tracking of voltage, current, temperature, and dust levels"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "Maintenance Management"
    p = tf.add_paragraph()
    p.text = "• Automated task scheduling based on system health"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• Priority-based alerts (Warning vs Critical)"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "Performance Analysis"
    p = tf.add_paragraph()
    p.text = "• Historical data visualization and trend analysis"
    p.level = 1

    # 7. Benefits & Impact
    slide = add_slide("Benefits & Impact")
    tf = slide.shapes.placeholders[1].text_frame
    p = tf.add_paragraph()
    p.text = "✅ Reduced Operational Costs"
    p.level = 0
    p = tf.add_paragraph()
    p.text = "   • Preventative maintenance is cheaper than emergency repairs"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "✅ Increased Energy Efficiency"
    p.level = 0
    p = tf.add_paragraph()
    p.text = "   • Optimal panel performance through timely cleaning"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "✅ Extended Equipment Life"
    p.level = 0
    p = tf.add_paragraph()
    p.text = "   • Early detection of degradation and faults"
    p.level = 1

    # 8. Future Scope
    slide = add_slide("Future Scope")
    tf = slide.shapes.placeholders[1].text_frame
    p = tf.add_paragraph()
    p.text = "• Mobile App Integration for remote monitoring"
    p = tf.add_paragraph()
    p.text = "• IoT Sensor Integration for direct data feed"
    p = tf.add_paragraph()
    p.text = "• Advanced Deep Learning models for better accuracy"
    p = tf.add_paragraph()
    p.text = "• Multi-site management capabilities"

    # 9. Conclusion
    slide = prs.slides.add_slide(prs.slide_layouts[0]) # Title Slide layout for end
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "Thank You"
    subtitle.text = "Questions?"

    # Save presentation
    output_path = r"C:\Users\Imamsab jamdar\OneDrive\Desktop\Predictive_maintenances-main\Predictive_Maintenance_Presentation.pptx"
    prs.save(output_path)
    print(f"Presentation saved to: {output_path}")

if __name__ == "__main__":
    create_presentation()
